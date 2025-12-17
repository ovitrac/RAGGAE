"""
RAGGAE/io/ingest.py

Unified document ingestion pipeline with multi-format support.

Provides:
- Auto-detection of file formats by extension
- Pandoc-first conversion (better quality) with native fallbacks
- ZIP/folder batch processing
- Consistent block extraction with provenance

Supported formats:
- PDF: PyMuPDF (native)
- DOCX/ODT: Pandoc → fallback python-docx/odfpy
- PPTX/ODP: Pandoc → fallback python-pptx/odfpy
- XLSX/ODS: Pandoc → fallback openpyxl/odfpy
- TXT/MD/CSV/JSON/XML: Native text loading
- ZIP: Extract and process recursively
- Folders: Walk and process all files

Author: Olivier Vitrac, PhD, HDR | olivier.vitrac@adservio.fr | Adservio | 2025-12-17
License: MIT
"""

from __future__ import annotations

import json
import logging
import os
import shutil
import subprocess
import tempfile
import zipfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Dict, List, Optional, Tuple, Union, Callable
from concurrent.futures import ThreadPoolExecutor, as_completed

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Data Structures
# ---------------------------------------------------------------------------

@dataclass
class TextBlock:
    """Generic text block with provenance."""
    text: str
    page: int = 1
    block: int = 0
    bbox: Tuple[float, float, float, float] = (0, 0, 0, 0)
    source_file: str = ""
    source_format: str = ""

    def as_metadata(self) -> Dict:
        return {
            "page": self.page,
            "block": self.block,
            "bbox": self.bbox,
            "file": self.source_file,
            "ext": self.source_format,
        }


@dataclass
class IngestResult:
    """Result of ingesting a single file."""
    file: str
    format: str
    blocks: List[TextBlock] = field(default_factory=list)
    status: str = "ok"  # ok, skipped, error
    error: Optional[str] = None
    conversion_method: str = ""  # pandoc, native, none

    @property
    def num_blocks(self) -> int:
        return len(self.blocks)


@dataclass
class BatchIngestResult:
    """Result of batch ingestion."""
    files_processed: int = 0
    files_skipped: int = 0
    files_errored: int = 0
    total_blocks: int = 0
    results: List[IngestResult] = field(default_factory=list)
    errors: List[str] = field(default_factory=list)


# ---------------------------------------------------------------------------
# Format Detection
# ---------------------------------------------------------------------------

# Extension to format mapping
EXTENSION_MAP = {
    # Documents
    "pdf": "pdf",
    "docx": "docx",
    "doc": "doc",
    "odt": "odt",
    "rtf": "rtf",
    # Presentations
    "pptx": "pptx",
    "ppt": "ppt",
    "odp": "odp",
    # Spreadsheets
    "xlsx": "xlsx",
    "xls": "xls",
    "ods": "ods",
    "csv": "csv",
    # Text
    "txt": "txt",
    "md": "md",
    "markdown": "md",
    "json": "json",
    "xml": "xml",
    "html": "html",
    "htm": "html",
    # Code (treat as text)
    "py": "txt",
    "java": "txt",
    "js": "txt",
    "ts": "txt",
    "c": "txt",
    "cpp": "txt",
    "h": "txt",
    "hpp": "txt",
    "rs": "txt",
    "go": "txt",
    "rb": "txt",
    "php": "txt",
    "sh": "txt",
    "bash": "txt",
    "yaml": "txt",
    "yml": "txt",
    "toml": "txt",
    "ini": "txt",
    "cfg": "txt",
    "conf": "txt",
    # Archives
    "zip": "zip",
}

# Formats that pandoc can convert
PANDOC_FORMATS = {"docx", "odt", "rtf", "pptx", "odp", "xlsx", "ods", "html", "csv"}

# Binary formats to skip
BINARY_FORMATS = {"jpg", "jpeg", "png", "gif", "bmp", "ico", "svg", "webp",
                  "mp3", "mp4", "wav", "avi", "mov", "mkv", "webm",
                  "exe", "dll", "so", "dylib", "bin", "dat",
                  "ttf", "otf", "woff", "woff2", "eot"}


def detect_format(path: Path) -> Optional[str]:
    """Detect file format from extension."""
    ext = path.suffix.lower().lstrip(".")
    if ext in BINARY_FORMATS:
        return None
    return EXTENSION_MAP.get(ext)


# ---------------------------------------------------------------------------
# Pandoc Integration
# ---------------------------------------------------------------------------

_pandoc_available: Optional[bool] = None


def check_pandoc() -> bool:
    """Check if pandoc is installed and accessible."""
    global _pandoc_available
    if _pandoc_available is None:
        try:
            result = subprocess.run(
                ["pandoc", "--version"],
                capture_output=True,
                timeout=5
            )
            _pandoc_available = result.returncode == 0
        except Exception:
            _pandoc_available = False
    return _pandoc_available


def convert_with_pandoc(
    input_path: Path,
    input_format: str,
    output_format: str = "plain"
) -> Optional[str]:
    """
    Convert file using pandoc.

    Args:
        input_path: Path to input file
        input_format: Pandoc input format (docx, odt, pptx, etc.)
        output_format: Pandoc output format (plain, markdown, etc.)

    Returns:
        Converted text or None on failure
    """
    if not check_pandoc():
        return None

    # Map our formats to pandoc format names
    pandoc_input_map = {
        "docx": "docx",
        "odt": "odt",
        "rtf": "rtf",
        "pptx": "pptx",
        "odp": "odp",
        "xlsx": "xlsx",  # pandoc 3.0+ supports xlsx
        "ods": "ods",
        "html": "html",
        "csv": "csv",
    }

    pandoc_format = pandoc_input_map.get(input_format)
    if not pandoc_format:
        return None

    try:
        result = subprocess.run(
            [
                "pandoc",
                "-f", pandoc_format,
                "-t", output_format,
                "--wrap=none",
                str(input_path)
            ],
            capture_output=True,
            text=True,
            timeout=60
        )
        if result.returncode == 0:
            return result.stdout
        else:
            logger.warning(f"Pandoc conversion failed: {result.stderr}")
            return None
    except subprocess.TimeoutExpired:
        logger.warning(f"Pandoc timeout for {input_path}")
        return None
    except Exception as e:
        logger.warning(f"Pandoc error: {e}")
        return None


# ---------------------------------------------------------------------------
# Native Loaders (Fallbacks)
# ---------------------------------------------------------------------------

def load_pdf_blocks(path: Path, min_chars: int = 20) -> List[TextBlock]:
    """Load PDF using PyMuPDF."""
    try:
        import fitz  # PyMuPDF
    except ImportError:
        raise ImportError("PyMuPDF not installed: pip install pymupdf")

    blocks = []
    doc = fitz.open(str(path))
    for page_num, page in enumerate(doc, 1):
        for block_num, block in enumerate(page.get_text("blocks")):
            text = block[4].strip() if len(block) > 4 else ""
            if len(text) >= min_chars:
                blocks.append(TextBlock(
                    text=text,
                    page=page_num,
                    block=block_num,
                    bbox=block[:4] if len(block) >= 4 else (0, 0, 0, 0),
                    source_file=path.name,
                    source_format="pdf"
                ))
    doc.close()
    return blocks


def load_docx_blocks(path: Path, min_chars: int = 20) -> List[TextBlock]:
    """Load DOCX using python-docx."""
    try:
        import docx
    except ImportError:
        raise ImportError("python-docx not installed: pip install python-docx")

    doc = docx.Document(str(path))
    blocks = []
    block_num = 0

    # Paragraphs
    for para in doc.paragraphs:
        text = (para.text or "").strip()
        if len(text) >= min_chars:
            blocks.append(TextBlock(
                text=text,
                page=1,
                block=block_num,
                source_file=path.name,
                source_format="docx"
            ))
            block_num += 1

    # Tables
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            text = " | ".join(c for c in cells if c)
            if len(text) >= min_chars:
                blocks.append(TextBlock(
                    text=text,
                    page=1,
                    block=block_num,
                    source_file=path.name,
                    source_format="docx"
                ))
                block_num += 1

    return blocks


def load_odt_blocks(path: Path, min_chars: int = 20) -> List[TextBlock]:
    """Load ODT using odfpy."""
    try:
        from odf import text, teletype
        from odf.opendocument import load
    except ImportError:
        raise ImportError("odfpy not installed: pip install odfpy")

    doc = load(str(path))
    paras = doc.getElementsByType(text.P)
    blocks = []

    for block_num, p in enumerate(paras):
        txt = (teletype.extractText(p) or "").strip()
        if len(txt) >= min_chars:
            blocks.append(TextBlock(
                text=txt,
                page=1,
                block=block_num,
                source_file=path.name,
                source_format="odt"
            ))

    return blocks


def load_pptx_blocks(path: Path, min_chars: int = 20) -> List[TextBlock]:
    """Load PPTX using python-pptx."""
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("python-pptx not installed: pip install python-pptx")

    prs = Presentation(str(path))
    blocks = []
    block_num = 0

    for slide_num, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if len(text) >= min_chars:
                    blocks.append(TextBlock(
                        text=text,
                        page=slide_num,
                        block=block_num,
                        source_file=path.name,
                        source_format="pptx"
                    ))
                    block_num += 1

            # Tables in slides
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    text = " | ".join(c for c in cells if c)
                    if len(text) >= min_chars:
                        blocks.append(TextBlock(
                            text=text,
                            page=slide_num,
                            block=block_num,
                            source_file=path.name,
                            source_format="pptx"
                        ))
                        block_num += 1

    return blocks


def load_xlsx_blocks(path: Path, min_chars: int = 20) -> List[TextBlock]:
    """Load XLSX using openpyxl."""
    try:
        from openpyxl import load_workbook
    except ImportError:
        raise ImportError("openpyxl not installed: pip install openpyxl")

    wb = load_workbook(str(path), read_only=True, data_only=True)
    blocks = []
    block_num = 0

    for sheet_num, sheet in enumerate(wb.worksheets, 1):
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c).strip() for c in row if c is not None]
            text = " | ".join(c for c in cells if c)
            if len(text) >= min_chars:
                blocks.append(TextBlock(
                    text=text,
                    page=sheet_num,
                    block=block_num,
                    source_file=path.name,
                    source_format="xlsx"
                ))
                block_num += 1

    wb.close()
    return blocks


def load_ods_blocks(path: Path, min_chars: int = 20) -> List[TextBlock]:
    """Load ODS using odfpy."""
    try:
        from odf.opendocument import load
        from odf import table, text, teletype
    except ImportError:
        raise ImportError("odfpy not installed: pip install odfpy")

    doc = load(str(path))
    blocks = []
    block_num = 0

    for sheet_num, sheet in enumerate(doc.getElementsByType(table.Table), 1):
        for row in sheet.getElementsByType(table.TableRow):
            cells = []
            for cell in row.getElementsByType(table.TableCell):
                cell_text = ""
                for p in cell.getElementsByType(text.P):
                    cell_text += teletype.extractText(p) or ""
                cells.append(cell_text.strip())

            text_row = " | ".join(c for c in cells if c)
            if len(text_row) >= min_chars:
                blocks.append(TextBlock(
                    text=text_row,
                    page=sheet_num,
                    block=block_num,
                    source_file=path.name,
                    source_format="ods"
                ))
                block_num += 1

    return blocks


def load_text_blocks(path: Path, min_chars: int = 20, fmt: str = "txt") -> List[TextBlock]:
    """Load plain text files."""
    try:
        content = path.read_text(encoding="utf-8", errors="ignore")
    except Exception as e:
        logger.warning(f"Failed to read {path}: {e}")
        return []

    blocks = []
    for block_num, para in enumerate(content.split("\n\n")):
        text = para.strip()
        if len(text) >= min_chars:
            blocks.append(TextBlock(
                text=text,
                page=1,
                block=block_num,
                source_file=path.name,
                source_format=fmt
            ))

    return blocks


def load_csv_blocks(path: Path, min_chars: int = 20) -> List[TextBlock]:
    """Load CSV files."""
    import csv

    blocks = []
    block_num = 0

    try:
        with open(path, newline="", encoding="utf-8", errors="ignore") as f:
            reader = csv.reader(f)
            for row in reader:
                text = " | ".join(c.strip() for c in row if c.strip())
                if len(text) >= min_chars:
                    blocks.append(TextBlock(
                        text=text,
                        page=1,
                        block=block_num,
                        source_file=path.name,
                        source_format="csv"
                    ))
                    block_num += 1
    except Exception as e:
        logger.warning(f"Failed to parse CSV {path}: {e}")

    return blocks


# ---------------------------------------------------------------------------
# Main Ingestion Functions
# ---------------------------------------------------------------------------

def ingest_file(
    path: Union[str, Path],
    min_chars: int = 20,
    prefer_pandoc: bool = True
) -> IngestResult:
    """
    Ingest a single file and extract text blocks.

    Args:
        path: Path to file
        min_chars: Minimum characters per block
        prefer_pandoc: Try pandoc first for supported formats

    Returns:
        IngestResult with extracted blocks or error
    """
    path = Path(path)

    if not path.exists():
        return IngestResult(
            file=str(path),
            format="unknown",
            status="error",
            error=f"File not found: {path}"
        )

    fmt = detect_format(path)
    if fmt is None:
        return IngestResult(
            file=str(path),
            format="binary",
            status="skipped",
            error="Binary or unsupported format"
        )

    result = IngestResult(file=str(path), format=fmt)

    try:
        blocks = []
        conversion_method = "native"

        # Try pandoc first for supported formats
        if prefer_pandoc and fmt in PANDOC_FORMATS:
            pandoc_text = convert_with_pandoc(path, fmt)
            if pandoc_text:
                conversion_method = "pandoc"
                # Split pandoc output into blocks
                for block_num, para in enumerate(pandoc_text.split("\n\n")):
                    text = para.strip()
                    if len(text) >= min_chars:
                        blocks.append(TextBlock(
                            text=text,
                            page=1,
                            block=block_num,
                            source_file=path.name,
                            source_format=fmt
                        ))

        # Fallback to native loaders if pandoc failed or not preferred
        if not blocks:
            conversion_method = "native"
            if fmt == "pdf":
                blocks = load_pdf_blocks(path, min_chars)
            elif fmt == "docx":
                blocks = load_docx_blocks(path, min_chars)
            elif fmt == "odt":
                blocks = load_odt_blocks(path, min_chars)
            elif fmt == "pptx":
                blocks = load_pptx_blocks(path, min_chars)
            elif fmt in ("xlsx", "xls"):
                blocks = load_xlsx_blocks(path, min_chars)
            elif fmt == "ods":
                blocks = load_ods_blocks(path, min_chars)
            elif fmt == "csv":
                blocks = load_csv_blocks(path, min_chars)
            elif fmt in ("txt", "md", "json", "xml", "html"):
                blocks = load_text_blocks(path, min_chars, fmt)
            else:
                # Try as plain text
                blocks = load_text_blocks(path, min_chars, fmt)

        result.blocks = blocks
        result.conversion_method = conversion_method
        result.status = "ok" if blocks else "skipped"
        if not blocks:
            result.error = "No text blocks extracted"

    except Exception as e:
        result.status = "error"
        result.error = str(e)
        logger.exception(f"Failed to ingest {path}")

    return result


def ingest_zip(
    zip_path: Union[str, Path],
    min_chars: int = 20,
    prefer_pandoc: bool = True,
    max_workers: int = 4
) -> BatchIngestResult:
    """
    Extract and ingest all files from a ZIP archive.

    Args:
        zip_path: Path to ZIP file
        min_chars: Minimum characters per block
        prefer_pandoc: Try pandoc first
        max_workers: Parallel workers for extraction

    Returns:
        BatchIngestResult with all extracted blocks
    """
    zip_path = Path(zip_path)
    batch = BatchIngestResult()

    with tempfile.TemporaryDirectory() as tmpdir:
        try:
            with zipfile.ZipFile(zip_path, "r") as zf:
                zf.extractall(tmpdir)
        except Exception as e:
            batch.errors.append(f"Failed to extract ZIP: {e}")
            return batch

        # Process all extracted files
        return ingest_folder(
            Path(tmpdir),
            min_chars=min_chars,
            prefer_pandoc=prefer_pandoc,
            max_workers=max_workers
        )


def ingest_folder(
    folder_path: Union[str, Path],
    min_chars: int = 20,
    prefer_pandoc: bool = True,
    max_workers: int = 4,
    recursive: bool = True,
    progress_callback: Optional[Callable[[int, int, str], None]] = None
) -> BatchIngestResult:
    """
    Ingest all files from a folder.

    Args:
        folder_path: Path to folder
        min_chars: Minimum characters per block
        prefer_pandoc: Try pandoc first
        max_workers: Parallel workers
        recursive: Process subfolders
        progress_callback: Called with (current, total, filename)

    Returns:
        BatchIngestResult with all extracted blocks
    """
    folder_path = Path(folder_path)
    batch = BatchIngestResult()

    if not folder_path.is_dir():
        batch.errors.append(f"Not a directory: {folder_path}")
        return batch

    # Collect all files
    if recursive:
        files = [f for f in folder_path.rglob("*") if f.is_file()]
    else:
        files = [f for f in folder_path.iterdir() if f.is_file()]

    # Filter out unsupported formats
    files = [f for f in files if detect_format(f) is not None]
    total = len(files)

    if max_workers > 1:
        # Parallel processing
        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            futures = {
                executor.submit(ingest_file, f, min_chars, prefer_pandoc): f
                for f in files
            }

            for i, future in enumerate(as_completed(futures)):
                f = futures[future]
                try:
                    result = future.result()
                    batch.results.append(result)

                    if result.status == "ok":
                        batch.files_processed += 1
                        batch.total_blocks += result.num_blocks
                    elif result.status == "skipped":
                        batch.files_skipped += 1
                    else:
                        batch.files_errored += 1
                        batch.errors.append(f"{f}: {result.error}")

                except Exception as e:
                    batch.files_errored += 1
                    batch.errors.append(f"{f}: {e}")

                if progress_callback:
                    progress_callback(i + 1, total, str(f.name))
    else:
        # Sequential processing
        for i, f in enumerate(files):
            result = ingest_file(f, min_chars, prefer_pandoc)
            batch.results.append(result)

            if result.status == "ok":
                batch.files_processed += 1
                batch.total_blocks += result.num_blocks
            elif result.status == "skipped":
                batch.files_skipped += 1
            else:
                batch.files_errored += 1
                batch.errors.append(f"{f}: {result.error}")

            if progress_callback:
                progress_callback(i + 1, total, str(f.name))

    return batch


def ingest_any(
    path: Union[str, Path],
    min_chars: int = 20,
    prefer_pandoc: bool = True,
    max_workers: int = 4,
    progress_callback: Optional[Callable[[int, int, str], None]] = None
) -> BatchIngestResult:
    """
    Universal ingestion: file, folder, or ZIP.

    Automatically detects input type and delegates to appropriate handler.

    Args:
        path: Path to file, folder, or ZIP
        min_chars: Minimum characters per block
        prefer_pandoc: Try pandoc first
        max_workers: Parallel workers for batch processing
        progress_callback: Progress reporting callback

    Returns:
        BatchIngestResult with all extracted blocks
    """
    path = Path(path)

    if not path.exists():
        batch = BatchIngestResult()
        batch.errors.append(f"Path not found: {path}")
        return batch

    if path.is_dir():
        return ingest_folder(
            path, min_chars, prefer_pandoc, max_workers, True, progress_callback
        )

    fmt = detect_format(path)
    if fmt == "zip":
        return ingest_zip(path, min_chars, prefer_pandoc, max_workers)

    # Single file
    result = ingest_file(path, min_chars, prefer_pandoc)
    batch = BatchIngestResult()
    batch.results.append(result)

    if result.status == "ok":
        batch.files_processed = 1
        batch.total_blocks = result.num_blocks
    elif result.status == "skipped":
        batch.files_skipped = 1
    else:
        batch.files_errored = 1
        batch.errors.append(f"{path}: {result.error}")

    return batch


# ---------------------------------------------------------------------------
# Utility Functions
# ---------------------------------------------------------------------------

def get_supported_extensions() -> List[str]:
    """Return list of supported file extensions."""
    return sorted(set(EXTENSION_MAP.keys()) - BINARY_FORMATS)


def is_pandoc_available() -> bool:
    """Check if pandoc is available."""
    return check_pandoc()


def get_format_info(ext: str) -> Dict[str, Any]:
    """Get information about a format."""
    fmt = EXTENSION_MAP.get(ext.lower().lstrip("."))
    return {
        "extension": ext,
        "format": fmt,
        "supported": fmt is not None,
        "pandoc_convertible": fmt in PANDOC_FORMATS,
        "binary": ext.lower().lstrip(".") in BINARY_FORMATS,
    }
